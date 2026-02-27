#!/usr/bin/env python3
"""Generate DCFC Investor Presentation PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # Dark navy background
BG_CARD = RGBColor(0x16, 0x21, 0x3A)       # Card background
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xE0)   # Primary accent
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)  # Success / positive
ACCENT_ORANGE = RGBColor(0xFF, 0x8F, 0x00) # Warning / attention
ACCENT_RED = RGBColor(0xE0, 0x40, 0x40)    # Critical
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xBC, 0xD0)
MID_GRAY = RGBColor(0x70, 0x80, 0x98)
BORDER_COLOR = RGBColor(0x2A, 0x3A, 0x55)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ── Helper Functions ──

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.02
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_paragraph(tf, text, size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2), font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet(tf, text, size=13, color=LIGHT_GRAY, level=0, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.level = level
    p.space_before = Pt(3)
    p.space_after = Pt(2)
    return p


def add_title_bar(slide, slide_num, title_text, subtitle_text=None):
    """Add slide number badge + title at top."""
    # Slide number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.35), Inches(0.55), Inches(0.40))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_BLUE
    badge.line.fill.background()
    badge.adjustments[0] = 0.15
    tf = badge.text_frame
    tf.word_wrap = False
    set_text(tf, str(slide_num), size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    tf.paragraphs[0].font.name = "Calibri"

    # Title
    tb = add_text_box(slide, Inches(1.3), Inches(0.25), Inches(10), Inches(0.55))
    set_text(tb.text_frame, title_text, size=26, color=WHITE, bold=True)

    # Subtitle
    if subtitle_text:
        tb2 = add_text_box(slide, Inches(1.3), Inches(0.75), Inches(10), Inches(0.35))
        set_text(tb2.text_frame, subtitle_text, size=13, color=MID_GRAY)

    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.1), Inches(12.1), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER_COLOR
    line.line.fill.background()


def make_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    return table


def style_table(table, header_color=ACCENT_BLUE, data=None, col_widths=None):
    """Style a table with header row and alternating data rows."""
    rows = len(table.rows)
    cols = len(table.columns)

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            # Fill
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.color.rgb = WHITE
                    p.font.bold = True
                    p.font.name = "Calibri"
            else:
                bg = BG_CARD if row_idx % 2 == 1 else RGBColor(0x12, 0x1C, 0x32)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.color.rgb = LIGHT_GRAY
                    p.font.name = "Calibri"
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)

    if data:
        for r, row_data in enumerate(data):
            for c, val in enumerate(row_data):
                table.cell(r, c).text = str(val)


def add_code_block(slide, left, top, width, height, text, font_size=9):
    shape = add_shape(slide, left, top, width, height, fill_color=RGBColor(0x0A, 0x10, 0x20), border_color=BORDER_COLOR)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(8)
    set_text(tf, text, size=font_size, color=RGBColor(0x80, 0xE0, 0x80), font_name="Courier New")
    return shape


def add_metric_card(slide, left, top, width, height, label, value, accent=ACCENT_BLUE):
    card = add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=BORDER_COLOR)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(8)
    set_text(tf, value, size=22, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, label, size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    return card


def add_bar(slide, left, top, width, height, fill_color=ACCENT_BLUE):
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = fill_color
    bar.line.fill.background()
    bar.adjustments[0] = 0.3
    return bar


# ══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)

# Large accent bar at top
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

# Title
tb = add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2))
set_text(tb.text_frame, "150 kW DC Fast Charger", size=44, color=WHITE, bold=True)
add_paragraph(tb.text_frame, "Modular, Open-Source, SiC-Based EV Charging Infrastructure", size=20, color=ACCENT_BLUE, bold=False, space_before=Pt(12))

# Feature cards
features = [
    ("In-House Power Electronics", "30 kW SiC modules (Rev 0.7)"),
    ("Open-Source Software", "EVerest (Linux Foundation Energy)"),
    ("400V & 800V Compatible", "150–1000 VDC output range"),
    ("Scalable Platform", "50 kW to 350 kW configurations"),
]
card_w = Inches(2.75)
card_h = Inches(1.1)
start_x = Inches(0.85)
y = Inches(4.0)
gap = Inches(0.2)

for i, (title, desc) in enumerate(features):
    x = start_x + i * (card_w + gap)
    card = add_shape(slide, x, y, card_w, card_h, fill_color=BG_CARD, border_color=BORDER_COLOR)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(10)
    set_text(tf, title, size=13, color=ACCENT_BLUE, bold=True)
    add_paragraph(tf, desc, size=11, color=LIGHT_GRAY)

# Footer
tb = add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5))
set_text(tb.text_frame, "Technical Investor Presentation  |  February 2026", size=11, color=MID_GRAY, alignment=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: Market Opportunity
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 2, "Market Opportunity", "Global DC Fast Charging Infrastructure Demand")

# Metric cards
metrics = [
    ("Global EV Sales (2025)", "~17M units", ACCENT_BLUE),
    ("DCFCs Needed by 2030", "~2.4M", ACCENT_GREEN),
    ("Avg Commercial Price (150kW)", "$30K–$60K", ACCENT_ORANGE),
    ("Our BOM @100 Units", "$21,500", ACCENT_GREEN),
]
for i, (label, value, color) in enumerate(metrics):
    add_metric_card(slide, Inches(0.6 + i * 3.15), Inches(1.4), Inches(2.95), Inches(1.1), label, value, color)

# Key insight box
shape = add_shape(slide, Inches(0.6), Inches(2.8), Inches(12.1), Inches(1.2), fill_color=RGBColor(0x0A, 0x1A, 0x30), border_color=ACCENT_BLUE)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(15)
tf.margin_top = Pt(12)
set_text(tf, "KEY INSIGHT", size=11, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "Most DCFC manufacturers use proprietary, vertically integrated software stacks. Our open-source EVerest approach eliminates software licensing costs ($2,000–$5,000/unit for competitors) and vendor lock-in while accelerating time-to-market.", size=13, color=LIGHT_GRAY)

# Market sizing table
data = [
    ["Segment", "Units/Year (est.)", "ASP", "TAM"],
    ["Highway Corridors", "200,000+", "$35K–$50K", "$7B–$10B"],
    ["Urban Charging Hubs", "300,000+", "$30K–$45K", "$9B–$13B"],
    ["Fleet / Depot Charging", "150,000+", "$35K–$60K", "$5B–$9B"],
    ["Destination / Retail", "100,000+", "$25K–$40K", "$2.5B–$4B"],
]
table = make_table(slide, 5, 4, Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.5))
style_table(table, data=data, col_widths=[Inches(3.5), Inches(2.5), Inches(3), Inches(3.1)])


# ══════════════════════════════════════════════════════════════
# SLIDE 3: Product Overview
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 3, "Product Overview", "150 kW CCS2 DC Fast Charger — Key Specifications")

data = [
    ["Parameter", "Specification"],
    ["AC Input", "3-phase 400–480 Vac, 50/60 Hz"],
    ["DC Output Voltage", "150–1000 VDC"],
    ["DC Output Current", "Up to 500 A"],
    ["Maximum Power", "150 kW (5× 30 kW modules)"],
    ["Peak Efficiency", ">96% (SiC-based)"],
    ["Power Factor", ">0.99"],
    ["Connector", "CCS Combo 2 (CCS1 optional)"],
    ["Charging Time (80%)", "400V: 20–30 min  ·  800V: 15–20 min"],
    ["Software", "EVerest (Linux Foundation Energy)"],
    ["Cooling", "Sealed cabinet + clip-on HVAC"],
    ["Enclosure", "IP55 target"],
    ["Standards", "IEC 61851, ISO 15118, OCPP 2.0.1"],
]
table = make_table(slide, 13, 2, Inches(0.6), Inches(1.4), Inches(7), Inches(5.5))
style_table(table, data=data, col_widths=[Inches(2.5), Inches(4.5)])

# Highlights on right
highlights = [
    ("EFFICIENCY", ">96%", "SiC MOSFETs throughout"),
    ("VOLTAGE RANGE", "150–1000V", "400V + 800V EV compatible"),
    ("SCALABLE", "50–350 kW", "Same platform, add modules"),
    ("UPTIME", ">98%", "N+1 redundancy, hot-swap"),
]
for i, (title, value, desc) in enumerate(highlights):
    y = Inches(1.5) + i * Inches(1.35)
    card = add_shape(slide, Inches(8.2), y, Inches(4.5), Inches(1.15), fill_color=BG_CARD, border_color=BORDER_COLOR)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    set_text(tf, title, size=10, color=MID_GRAY, bold=True)
    add_paragraph(tf, value, size=24, color=ACCENT_BLUE, bold=True, space_before=Pt(2))
    add_paragraph(tf, desc, size=11, color=LIGHT_GRAY, space_before=Pt(0))


# ══════════════════════════════════════════════════════════════
# SLIDE 4: System Architecture
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 4, "System Architecture", "4-Level Hierarchical Control with Domain-Isolated CAN Buses")

arch_text = (
    "LEVEL 3: CLOUD / BACKEND\n"
    "  CSMS (OCPP 2.0.1)  ·  Fleet Management  ·  Billing\n"
    "           │ WebSocket / 4G / Ethernet\n"
    "LEVEL 2: STATION CONTROLLER\n"
    "  Phytec phyCORE-AM62x Main Controller (Linux + EVerest)\n"
    "  EvseManager · EnergyMgr · OCPP · Auth · ISO 15118\n"
    "  SafetySupervisorBSP · PowerModuleDriver · HvacDriver\n"
    "     │ CAN #2        │ CAN #1 (FD)      │ CAN #3\n"
    "LEVEL 1: MODULE CONTROLLERS\n"
    "  Safety Supv.    PDU-Micro Master      HVAC Unit\n"
    "  (STM32, SIL2)   (Module 0 PIM)       (STM32/RP2350)\n"
    "                   │  │  │  │\n"
    "                  M1  M2  M3  M4  (30kW each)\n"
    "LEVEL 0: HARDWARE\n"
    "  Gate Drivers · Sensors · Contactors · Actuators"
)
add_code_block(slide, Inches(0.6), Inches(1.4), Inches(7.5), Inches(5.0), arch_text, font_size=11)

# Key points on right
tb = add_text_box(slide, Inches(8.5), Inches(1.4), Inches(4.2), Inches(5.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Design Principles", size=16, color=WHITE, bold=True)
bullets = [
    "Safety-critical functions on dedicated STM32 (SIL 2) — never on Linux",
    "Three independent CAN buses isolate safety, power, and thermal domains",
    "Hardware interlock chain operates independently of all software",
    "EVerest microservices communicate via MQTT — modular and testable",
    "Main controller monitors safety but does NOT own the safety loop",
    "PDU-Micro modules self-coordinate via droop-based current sharing",
]
for b in bullets:
    add_bullet(tf, b, size=12, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: Modular Power Architecture
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 5, "Modular Power Architecture", "5× 30 kW Hot-Swappable Modules — Scalable to 350 kW")

# Configuration table
data = [
    ["Configuration", "Modules", "Power", "BOM @100", "BOM @500"],
    ["Single CCS2", "5× 30 kW", "150 kW", "$21,503", "$19,486"],
    ["Dual CCS2 (shared)", "5× 30 kW", "150 kW", "~$25,700", "~$23,700"],
    ["Single CCS2", "10× 30 kW", "300 kW", "~$32,900", "~$28,900"],
    ["Dual CCS2", "12× 30 kW", "350 kW", "~$36,900", "~$32,900"],
]
table = make_table(slide, 5, 5, Inches(0.6), Inches(1.4), Inches(12.1), Inches(2.2))
style_table(table, data=data, col_widths=[Inches(3), Inches(2), Inches(2), Inches(2.5), Inches(2.6)])

# Benefits cards
benefits = [
    ("Hot-Swappable", "Replace modules without powering down the entire system. Reduces MTTR to minutes."),
    ("N+1 Redundancy", "Optional extra module ensures continued operation during single-module failure."),
    ("Graceful Degradation", "Module failure reduces power proportionally — charger stays operational at reduced output."),
    ("Field-Scalable", "Start at 150 kW, upgrade to 300 kW by adding modules. Same cabinet, same software."),
]
for i, (title, desc) in enumerate(benefits):
    x = Inches(0.6) + i * Inches(3.1)
    card = add_shape(slide, x, Inches(4.0), Inches(2.9), Inches(2.0), fill_color=BG_CARD, border_color=BORDER_COLOR)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(10)
    set_text(tf, title, size=14, color=ACCENT_BLUE, bold=True)
    add_paragraph(tf, desc, size=11, color=LIGHT_GRAY, space_before=Pt(8))


# ══════════════════════════════════════════════════════════════
# SLIDE 6: PDU-Micro Power Module
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 6, "Power Module — PDU-Micro (30 kW)", "In-House Designed SiC Power Converter — Core IP")

data = [
    ["Parameter", "Specification"],
    ["Rated Power", "30 kW continuous"],
    ["PFC Topology", "Vienna Rectifier, 140 kHz"],
    ["DC-DC Topology", "3-phase Interleaved DAB"],
    ["SiC Devices (PFC)", "MSCSM120VR1M16CTPAG module"],
    ["SiC Devices (DAB)", "MSC080SMA120B discrete"],
    ["Output Voltage", "150–1000 VDC"],
    ["Output Current", "100 A per module"],
    ["DC Bus", "700–920 VDC (adaptive)"],
    ["Controller", "Dual dsPIC33CH512MP506"],
    ["Communication", "CAN-FD 500k/2M bps"],
    ["Peak Efficiency", "96.3% (system-level)"],
    ["Status", "Rev 0.7, firmware in development"],
]
table = make_table(slide, 13, 2, Inches(0.6), Inches(1.4), Inches(6.5), Inches(5.0))
style_table(table, data=data, col_widths=[Inches(2.5), Inches(4.0)])

# Cost metrics on right
cost_metrics = [
    ("Unit Cost @100", "$1,806", ACCENT_BLUE),
    ("Unit Cost @500", "$1,436", ACCENT_GREEN),
    ("Cost per kW @100", "$60/kW", ACCENT_BLUE),
    ("Cost per kW @500", "$48/kW", ACCENT_GREEN),
]
for i, (label, value, color) in enumerate(cost_metrics):
    add_metric_card(slide, Inches(7.8), Inches(1.4 + i * 1.25), Inches(4.9), Inches(1.05), label, value, color)

# Cost drivers
tb = add_text_box(slide, Inches(7.8), Inches(5.6), Inches(4.9), Inches(1.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Top Cost Drivers:", size=12, color=WHITE, bold=True)
add_bullet(tf, "DAB SiC MOSFETs — 23%", size=11, color=LIGHT_GRAY)
add_bullet(tf, "Vienna PFC module — 13.3%", size=11, color=LIGHT_GRAY)
add_bullet(tf, "PCB fabrication — 13%", size=11, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: Safety Architecture
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 7, "Safety Architecture", "Dual-Layer: Hardware Interlock Chain + SIL 2 Software Supervision")

# Safety chain diagram
chain_text = (
    "+24V Safety Supply\n"
    "    │\n"
    "    ▼\n"
    "E-STOP (NC) ─── Door Interlock (NC) ─── IMD Contact (NC)\n"
    "    │                                         │\n"
    "RCD Contact (NC) ─── Thermal Trip (NC) ─── Safety Relay\n"
    "    │\n"
    "    ▼  OUTPUT:\n"
    "AC Contactor  ·  DC Contactor  ·  Pre-charge Relay\n"
    "\n"
    "ANY BREAK IN CHAIN = IMMEDIATE SHUTDOWN\n"
    "Main controller/EVerest MONITORS but does NOT own the safety loop"
)
add_code_block(slide, Inches(0.6), Inches(1.4), Inches(6.5), Inches(3.2), chain_text, font_size=10)

# Protection table
data = [
    ["Protection", "Response", "Method"],
    ["Over-Voltage (OVP)", "< 1 µs", "Hardware comparator"],
    ["Over-Current (OCP)", "< 1 ms", "Safety supervisor ADC"],
    ["Ground Fault", "< 30 ms", "IMD relay (Bender)"],
    ["Insulation Failure", "CableCheck", "IMD measurement"],
    ["Main Controller Heartbeat Loss", "2 s timeout", "Orderly shutdown"],
    ["Contactor Weld", "Post-open", "Aux NC + voltage verify"],
    ["E-Stop", "Immediate", "Hardware loop break"],
]
table = make_table(slide, 8, 3, Inches(7.5), Inches(1.4), Inches(5.2), Inches(3.2))
style_table(table, data=data, col_widths=[Inches(2.0), Inches(1.2), Inches(2.0)])

# Safety state machine
states_text = "States: INIT → IDLE → AC_CLOSE → PRECHARGE → DC_CLOSE → CHARGING → SHUTDOWN → FAULT"
shape = add_shape(slide, Inches(0.6), Inches(4.9), Inches(12.1), Inches(0.6), fill_color=RGBColor(0x0A, 0x1A, 0x30), border_color=ACCENT_BLUE)
tf = shape.text_frame
tf.margin_left = Pt(12)
tf.margin_top = Pt(6)
set_text(tf, states_text, size=12, color=ACCENT_BLUE, bold=True)

# Target badge
tb = add_text_box(slide, Inches(0.6), Inches(5.8), Inches(12.1), Inches(0.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Certification Target: IEC 61508 SIL 2 / ISO 13849 PLd", size=14, color=ACCENT_GREEN, bold=True)
add_paragraph(tf, "Dedicated STM32 bare-metal controller · MISRA C:2012 · >90% branch coverage · FMEA documentation", size=12, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 8: Software Stack — EVerest
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 8, "Software Stack — EVerest Framework", "Open-Source Charging OS — Linux Foundation Energy")

# Stack diagram
stack_text = (
    "┌───────────────────────────────────────────┐\n"
    "│          APPLICATION LAYER                 │\n"
    "│  EvseManager · EnergyManager · Auth · OCPP │\n"
    "├───────────────────────────────────────────┤\n"
    "│          MIDDLEWARE LAYER                  │\n"
    "│  ISO 15118 · OCPP 2.0.1 · IEC 61851      │\n"
    "│  Safety State Machine                     │\n"
    "├───────────────────────────────────────────┤\n"
    "│     HARDWARE ABSTRACTION LAYER            │\n"
    "│  CAN · Ethernet · GPIO · UART · SPI       │\n"
    "├───────────────────────────────────────────┤\n"
    "│      LINUX OS (phyCORE-AM62x: 4-core A53, 4GB)      │\n"
    "└───────────────────────────────────────────┘"
)
add_code_block(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(3.5), stack_text, font_size=10)

# Why EVerest
tb = add_text_box(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(3.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Why EVerest", size=16, color=WHITE, bold=True)
reasons = [
    "$0 licensing — Apache 2.0 open source",
    "Pre-built ISO 15118, OCPP, IEC 61851 stacks",
    "Microservices via MQTT — independent processes",
    "Linux Foundation Energy backing — major OEM support",
    "Conformance-tested protocol implementations",
    "Active community with regular releases",
]
for r in reasons:
    add_bullet(tf, r, size=12, color=LIGHT_GRAY)

# Custom modules section
tb = add_text_box(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(0.4))
set_text(tb.text_frame, "Custom EVerest Modules (Our IP)", size=14, color=WHITE, bold=True)

modules = [
    ("SafetySupervisorBSP", "CAN #2 bridge to STM32\nevse_board_support, isolation_monitor, ac_rcd, connector_lock"),
    ("PowerModuleDriver", "CAN-FD bridge to PDU-Micro master\npower_supply_DC, powermeter, module shedding"),
    ("HvacDriver", "CAN #3 bridge to HVAC unit\nthermal_management, derating logic"),
]
for i, (name, desc) in enumerate(modules):
    x = Inches(0.6) + i * Inches(4.1)
    card = add_shape(slide, x, Inches(5.7), Inches(3.9), Inches(1.2), fill_color=BG_CARD, border_color=ACCENT_BLUE)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(8)
    set_text(tf, name, size=12, color=ACCENT_BLUE, bold=True, font_name="Courier New")
    add_paragraph(tf, desc, size=10, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 9: Communication Protocols
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 9, "Communication & Protocol Stack", "Standards-Compliant Vehicle and Network Communication")

# Vehicle side table
tb = add_text_box(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.4))
set_text(tb.text_frame, "Vehicle Side (EV ↔ Charger)", size=14, color=ACCENT_BLUE, bold=True)

data = [
    ["Protocol", "Standard", "Function"],
    ["Control Pilot", "IEC 61851-1", "Basic signaling (States A–F)"],
    ["PLC (HomePlug GreenPHY)", "ISO 15118-3", "Physical layer for HLC"],
    ["DIN SPEC 70121", "DIN 70121", "Legacy DC fast charging"],
    ["ISO 15118-2", "ISO 15118-2", "Plug & Charge, smart charging"],
]
table = make_table(slide, 5, 3, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.2))
style_table(table, data=data, col_widths=[Inches(2.2), Inches(1.5), Inches(2.1)])

# Network side table
tb = add_text_box(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(0.4))
set_text(tb.text_frame, "Network Side (Charger ↔ Backend)", size=14, color=ACCENT_BLUE, bold=True)

data = [
    ["Protocol", "Version", "Function"],
    ["OCPP", "2.0.1", "Remote monitoring, billing"],
    ["TLS", "1.2/1.3", "Secure Plug & Charge certs"],
    ["MQTT", "Internal", "EVerest inter-module messaging"],
    ["HTTP/REST", "Service", "Diagnostic web UI"],
]
table = make_table(slide, 5, 3, Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.2))
style_table(table, data=data, col_widths=[Inches(1.8), Inches(1.5), Inches(2.5)])

# Internal CAN buses
tb = add_text_box(slide, Inches(0.6), Inches(4.2), Inches(12.1), Inches(0.4))
set_text(tb.text_frame, "Internal Communication Buses", size=14, color=ACCENT_BLUE, bold=True)

buses = [
    ("CAN #1 (CAN-FD)", "500k/2M bps", "Main Controller ↔ PDU-Micro modules", "Power control & telemetry"),
    ("CAN #2", "500 kbps", "Main Controller ↔ Safety Supervisor (STM32)", "Safety state, contactor control"),
    ("CAN #3", "250 kbps", "Main Controller ↔ HVAC Unit", "Thermal management"),
]
data_table = [["Bus", "Speed", "Path", "Function"]] + [[b[0], b[1], b[2], b[3]] for b in buses]
table = make_table(slide, 4, 4, Inches(0.6), Inches(4.6), Inches(12.1), Inches(1.6))
style_table(table, data=data_table, col_widths=[Inches(2.5), Inches(1.8), Inches(4.3), Inches(3.5)])

# Connectivity
tb = add_text_box(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5))
set_text(tb.text_frame, "Connectivity:  Ethernet (LAN)  ·  5G Modem (WAN)  ·  WiFi (Service)  ·  PoE Switch  ·  PLC Modem (ISO 15118)", size=12, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 10: Thermal Management
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 10, "Thermal Management — Clip-On HVAC", "Sealed Cabinet with External Heat Rejection")

hvac_text = (
    "  HVAC CLIP-ON UNIT              MAIN CABINET\n"
    "┌─────────────────┐           ┌─────────────────┐\n"
    "│ Condenser +     │           │ Power Modules    │\n"
    "│ Compressor      │           │ M1  M2  M3  M4   │\n"
    "│      │          │           │ M5                │\n"
    "│ Evaporator      │ Cold Air  │                   │\n"
    "│ Coil       ══════════════>  │ Internal          │\n"
    "│      │          │           │ Heatsinks         │\n"
    "│ Blower Fans     │ Hot Air   │                   │\n"
    "│           <══════════════   │ Hot Air Return    │\n"
    "│                 │           │                   │\n"
    "│ HVAC Controller │◄─CAN#3──►│ Main ECU (AM62x)      │\n"
    "└─────────────────┘           └─────────────────┘\n"
    "\n"
    "Internal air NEVER mixes with ambient → IP55 sealed"
)
add_code_block(slide, Inches(0.6), Inches(1.4), Inches(6.5), Inches(3.8), hvac_text, font_size=10)

# Specs table
data = [
    ["Parameter", "Specification"],
    ["Cooling Capacity", "15–25 kW thermal"],
    ["Airflow", "2,000–4,000 CFM (variable)"],
    ["Operating Range", "-20°C to +50°C ambient"],
    ["Refrigerant", "R-410A or R-32"],
    ["Mounting", "Side, rear, or top clip-on"],
    ["Interface", "CAN bus + 24V power"],
    ["Serviceability", "Field-replaceable, single unit"],
]
table = make_table(slide, 8, 2, Inches(7.5), Inches(1.4), Inches(5.2), Inches(3.4))
style_table(table, data=data, col_widths=[Inches(2.2), Inches(3.0)])

# Derating chain
shape = add_shape(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(0.9), fill_color=RGBColor(0x0A, 0x1A, 0x30), border_color=ACCENT_ORANGE)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(15)
tf.margin_top = Pt(8)
set_text(tf, "THERMAL DERATING CHAIN", size=12, color=ACCENT_ORANGE, bold=True)
add_paragraph(tf, "Cabinet temp rising → HVAC signals via CAN #3 → EnergyManager reduces output power → Prevents thermal shutdown", size=12, color=LIGHT_GRAY)
add_paragraph(tf, "Modes: Idle → Pre-Cool → Active Cooling → Max Cooling → Derating", size=11, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 11: Competitive Analysis
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 11, "Competitive Analysis", "150 kW Class DC Fast Chargers")

data = [
    ["Feature", "Our Design", "ABB Terra 184", "Tritium PKM150", "Delta UFC 150", "Kempower S"],
    ["Max Power", "150 kW (to 350)", "180 kW", "150 kW", "150 kW", "40–600 kW"],
    ["Efficiency", ">96% (SiC)", "95%", "~95%", ">94%", "N/A"],
    ["DC Output", "150–1000V", "150–920V", "150–920V", "200–1000V", "50–920V"],
    ["Max Current", "500 A", "200 A", "350 A", "500 A", "500 A"],
    ["IP Rating", "IP55 (target)", "IP54", "IP65", "IP55", "IP54"],
    ["Software", "EVerest (open)", "Proprietary", "Proprietary", "Proprietary", "Proprietary"],
    ["SiC MOSFETs", "Yes", "Unconfirmed", "Unconfirmed", "Unconfirmed", "Unconfirmed"],
    ["OCPP", "2.0.1", "1.6J", "1.6J / 2.0.1", "1.5S / 1.6J", "1.6J / 2.0.1"],
    ["Est. Retail", "$35K–40K", "$40K–50K", "$35K–45K", "$30K–45K", "$35K–50K"],
]
table = make_table(slide, 10, 6, Inches(0.6), Inches(1.4), Inches(12.1), Inches(4.0))
style_table(table, data=data, col_widths=[Inches(1.5), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.2)])

# Highlight our column
for row_idx in range(1, 10):
    cell = table.cell(row_idx, 1)
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = ACCENT_GREEN
        p.font.bold = True

# Differentiators
tb = add_text_box(slide, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Our Differentiators:", size=14, color=WHITE, bold=True)
diffs = [
    "Open-source software — no vendor lock-in, no licensing fees",
    "SiC throughout — highest efficiency in class",
    "Full 1000V output — future-proofed for 800V EV wave",
    "Modular power — same platform scales 50–350 kW",
    "In-house power electronics — full control of cost and supply chain",
]
for d in diffs:
    add_bullet(tf, "» " + d, size=11, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 12: BOM Breakdown
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 12, "Bill of Materials Breakdown", "150 kW Single CCS2 — Component Cost @100 Units")

# BOM bar chart (visual bars)
bom_items = [
    ("Power Modules (5×30 kW)", 10087, 46.9),
    ("CCS Connector + Cable", 2180, 10.1),
    ("Cabinet and Enclosure", 2145, 10.0),
    ("DC Output Protection", 1785, 8.3),
    ("HVAC Clip-On Unit", 1720, 8.0),
    ("AC Input and Protection", 1496, 7.0),
    ("Control Electronics", 730, 3.4),
    ("Cooling System", 635, 3.0),
    ("Networking + Comms", 375, 1.7),
    ("User Interface", 200, 0.9),
    ("Auxiliary PSU", 150, 0.7),
]

max_bar_width = Inches(5.5)
bar_h = Inches(0.28)
y_start = Inches(1.5)
label_x = Inches(0.6)
bar_x = Inches(3.5)
val_x = Inches(9.3)

for i, (name, cost, pct) in enumerate(bom_items):
    y = y_start + i * Inches(0.42)
    # Label
    tb = add_text_box(slide, label_x, y, Inches(2.8), bar_h)
    set_text(tb.text_frame, name, size=10, color=LIGHT_GRAY)
    # Bar
    bar_w = int(max_bar_width * (pct / 46.9))
    color = ACCENT_BLUE if i == 0 else RGBColor(0x00, 0x70, 0xAA) if pct > 7 else MID_GRAY
    add_bar(slide, bar_x, y + Pt(2), bar_w, Inches(0.22), fill_color=color)
    # Value
    tb = add_text_box(slide, val_x, y, Inches(3), bar_h)
    set_text(tb.text_frame, f"${cost:,}  ({pct}%)", size=10, color=LIGHT_GRAY)

# Total
y_total = y_start + 11 * Inches(0.42)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x, y_total, max_bar_width, Pt(1))
line.fill.solid()
line.fill.fore_color.rgb = BORDER_COLOR
line.line.fill.background()

tb = add_text_box(slide, label_x, y_total + Pt(4), Inches(2.8), Inches(0.3))
set_text(tb.text_frame, "TOTAL BOM", size=12, color=WHITE, bold=True)
tb = add_text_box(slide, val_x, y_total + Pt(4), Inches(3), Inches(0.3))
set_text(tb.text_frame, "$21,503", size=14, color=ACCENT_GREEN, bold=True)

# Margin table
data = [
    ["Volume", "BOM Cost", "Target Retail", "Gross Margin"],
    ["@100 units", "$21,503", "$35,000–40,000", "40–46%"],
    ["@500 units", "$19,486", "$30,000–35,000", "35–44%"],
]
table = make_table(slide, 3, 4, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.9))
style_table(table, data=data, col_widths=[Inches(2.5), Inches(3), Inches(3.3), Inches(3.3)])


# ══════════════════════════════════════════════════════════════
# SLIDE 13: Cost Optimization
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 13, "Cost Optimization Roadmap", "Path to Further BOM Reduction")

data = [
    ["Initiative", "Savings/Unit", "Timeline"],
    ["Chinese SiC alternatives (DAB MOSFETs)", "$500–$650", "Near-term"],
    ["Volume pricing on COTS (CCS, contactors)", "$1,000–$1,500", "@500+ units"],
    ["PCB panel optimization + assembly volume", "$300–$500", "@500+ units"],
    ["Backplane redesign (simplified busbar)", "$100–$200", "Rev B"],
    ["In-house HVAC manufacturing", "$300–$500", "Year 2"],
]
table = make_table(slide, 6, 3, Inches(0.6), Inches(1.4), Inches(12.1), Inches(2.8))
style_table(table, data=data, col_widths=[Inches(5.5), Inches(3.3), Inches(3.3)])

# Projection cards
projections = [
    ("BOM @100 (Current)", "$21,503", ACCENT_BLUE),
    ("BOM @500", "$19,486", ACCENT_BLUE),
    ("Projected BOM @1000", "$16,500–$17,500", ACCENT_GREEN),
    ("Gross Margin @1000", "~47%", ACCENT_GREEN),
]
for i, (label, value, color) in enumerate(projections):
    add_metric_card(slide, Inches(0.6 + i * 3.15), Inches(4.5), Inches(2.95), Inches(1.1), label, value, color)

# Software advantage
shape = add_shape(slide, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.0), fill_color=RGBColor(0x0A, 0x1A, 0x30), border_color=ACCENT_GREEN)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(15)
tf.margin_top = Pt(10)
set_text(tf, "SOFTWARE COST ADVANTAGE", size=12, color=ACCENT_GREEN, bold=True)
add_paragraph(tf, "Competitors spend $2,000–$5,000 per unit on proprietary software licensing. Our EVerest-based stack has $0 per-unit software cost — a structural margin advantage at every volume.", size=12, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 14: Standards & Certification
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 14, "Standards & Certification Plan", "International Compliance Matrix")

data = [
    ["Standard", "Scope", "Status"],
    ["IEC 61851-1", "General EV charging requirements", "Design phase"],
    ["IEC 61851-23", "DC charging station requirements", "Design phase"],
    ["IEC 62368-1", "Electrical safety", "Testing planned"],
    ["IEC 61000-6-2/6-4", "EMC (emissions + immunity)", "Pre-compliance M11"],
    ["ISO 15118-2/3", "Vehicle-to-grid communication", "EVerest pre-certified"],
    ["ISO 13849 (PLd)", "Safety of machinery", "Safety supervisor design"],
    ["IEC 61508 (SIL 2)", "Functional safety", "Safety supervisor design"],
    ["OCPP 2.0.1", "Charge point protocol", "EVerest (libocpp) certified"],
    ["UL 2202", "US market certification", "Planned for US launch"],
]
table = make_table(slide, 10, 3, Inches(0.6), Inches(1.4), Inches(7.5), Inches(4.2))
style_table(table, data=data, col_widths=[Inches(2.0), Inches(3.2), Inches(2.3)])

# Regional certifications
tb = add_text_box(slide, Inches(8.5), Inches(1.4), Inches(4.2), Inches(0.4))
set_text(tb.text_frame, "Certification Bodies", size=14, color=ACCENT_BLUE, bold=True)

regions = [
    ("Europe", "TÜV", "CE, IEC 61851, EMC"),
    ("North America", "UL", "UL 2202, NEC Art. 625"),
    ("India", "BIS", "IS 17017 (planned)"),
    ("China", "CQC", "GB/T standards (future)"),
]
for i, (region, body, cert) in enumerate(regions):
    card = add_shape(slide, Inches(8.5), Inches(1.9 + i * Inches(1.0)), Inches(4.2), Inches(0.85), fill_color=BG_CARD, border_color=BORDER_COLOR)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(6)
    set_text(tf, region, size=12, color=WHITE, bold=True)
    add_paragraph(tf, f"{body} — {cert}", size=11, color=LIGHT_GRAY)

# Testing phases
tb = add_text_box(slide, Inches(8.5), Inches(5.8), Inches(4.2), Inches(1.2))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Testing Phases:", size=13, color=WHITE, bold=True)
phases = ["Type testing (electrical, EMC, safety)", "Protocol conformance (CharIN, OCA)", "Field testing (installation verification)", "SIL 2 / PLd assessment"]
for ph in phases:
    add_bullet(tf, ph, size=11, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 15: Development Timeline
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 15, "Development Plan — 14-Month Timeline", "8 Parallel Workstreams to Certification-Ready Prototype")

# Gantt chart
workstreams = [
    ("WS-1", "PDU-Micro Modules", 1, 10, ACCENT_BLUE),
    ("WS-2", "CAN Protocol Bridge", 2, 4, RGBColor(0x00, 0x70, 0xAA)),
    ("WS-3", "Safety Supervisor FW", 1, 6, ACCENT_RED),
    ("WS-4", "EVerest Integration", 2, 10, ACCENT_GREEN),
    ("WS-5", "Cabinet Hardware", 1, 6, ACCENT_ORANGE),
    ("WS-6", "HVAC Controller FW", 3, 7, RGBColor(0x80, 0x60, 0xC0)),
    ("WS-7", "System Integration", 7, 10, RGBColor(0x00, 0xA0, 0x80)),
    ("WS-8", "Certification", 9, 14, RGBColor(0xC0, 0x80, 0x00)),
]

gantt_left = Inches(3.2)
gantt_width = Inches(9.5)
month_width = gantt_width / 14
y_start = Inches(1.6)
row_h = Inches(0.55)

# Month headers
for m in range(1, 15):
    x = gantt_left + month_width * (m - 1)
    tb = add_text_box(slide, x, Inches(1.25), month_width, Inches(0.3))
    set_text(tb.text_frame, f"M{m}", size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Workstream rows
for i, (ws_id, ws_name, start, end, color) in enumerate(workstreams):
    y = y_start + i * row_h
    # Label
    tb = add_text_box(slide, Inches(0.3), y + Pt(4), Inches(2.8), Inches(0.3))
    tf = tb.text_frame
    set_text(tf, f"{ws_id}: ", size=10, color=color, bold=True)
    run = tf.paragraphs[0].add_run()
    run.text = ws_name
    run.font.size = Pt(10)
    run.font.color.rgb = LIGHT_GRAY
    run.font.name = "Calibri"
    # Bar
    bar_left = gantt_left + month_width * (start - 1)
    bar_width = month_width * (end - start + 1)
    add_bar(slide, int(bar_left), y + Pt(4), int(bar_width), Inches(0.30), fill_color=color)

# Critical path note
y_note = y_start + 8 * row_h + Pt(10)
shape = add_shape(slide, Inches(0.6), y_note, Inches(12.1), Inches(0.7), fill_color=RGBColor(0x0A, 0x1A, 0x30), border_color=ACCENT_RED)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(12)
tf.margin_top = Pt(8)
set_text(tf, "CRITICAL PATH:", size=11, color=ACCENT_RED, bold=True)
run = tf.paragraphs[0].add_run()
run.text = "  Module Firmware → CAN Bridge → EVerest Modules → System Integration → Certification"
run.font.size = Pt(11)
run.font.color.rgb = LIGHT_GRAY
run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════
# SLIDE 16: Key Milestones
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 16, "Key Milestones", "14-Month Execution Plan — Monthly Gate Criteria")

data = [
    ["Month", "Milestone", "Gate Criteria"],
    ["M2", "Safety supervisor board powered", "CAN #2 loopback verified"],
    ["M3", "First PDU-Micro module producing DC", "30 kW output on bench"],
    ["M4", "phyCORE-AM62x running EVerest", "SafetySupervisorBSP ↔ STM32"],
    ["M5", "PowerModuleDriver controlling 1 module", "CAN-FD setpoint control"],
    ["M6", "Cabinet assembled, 5 modules installed", "Backplane wired, mech fit"],
    ["M7", "HVAC clip-on operational", "Thermal derating chain verified"],
    ["M8", "First complete charging session", "ISO 15118 w/ EV simulator"],
    ["M9", "Full-power 150 kW sustained test", "Thermal steady-state OK"],
    ["M10", "OCPP backend + HMI complete", "Remote monitoring live"],
    ["M11", "EMC pre-compliance pass", "EN 55032 Class B"],
    ["M12–14", "Certification testing complete", "IEC 61851, 62368, EMC"],
]
table = make_table(slide, 12, 3, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.5))
style_table(table, data=data, col_widths=[Inches(1.3), Inches(5.4), Inches(5.4)])

# Highlight key milestones
for row_idx in [8, 9, 11]:  # M8, M9, M12-14
    for col_idx in range(3):
        cell = table.cell(row_idx, col_idx)
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = ACCENT_GREEN
            p.font.bold = True


# ══════════════════════════════════════════════════════════════
# SLIDE 17: Team & Resources
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 17, "Team & Resource Requirements", "8–10 Engineers (3–4 Shared from Existing PDU-Micro Team)")

data = [
    ["Role", "Count", "Key Skills"],
    ["Power Electronics Engineer", "1–2", "PDU-Micro hardware, magnetics, thermal"],
    ["Embedded FW Engineer (Safety)", "1", "STM32, bare-metal C, MISRA C, IEC 61508"],
    ["Embedded FW Engineer (Controls)", "1–2", "dsPIC33CH, power converter control loops"],
    ["Linux/EVerest Software Engineer", "1–2", "C++, Linux, EVerest, SocketCAN, MQTT"],
    ["Protocol Engineer", "1", "ISO 15118, OCPP 2.0.1, CAN bus, PLC"],
    ["Electrical/Mechanical Engineer", "1", "Cabinet design, busbar, wiring, IP55"],
    ["HVAC Engineer", "0.5", "Refrigeration, compressor (part-time)"],
    ["QA / Certification Engineer", "1", "IEC 61851, IEC 62368, EMC, test planning"],
]
table = make_table(slide, 9, 3, Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.8))
style_table(table, data=data, col_widths=[Inches(3.5), Inches(1.3), Inches(7.3)])

# Leverage notes
shape = add_shape(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.3), fill_color=BG_CARD, border_color=BORDER_COLOR)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(15)
tf.margin_top = Pt(10)
set_text(tf, "Team Leverage", size=14, color=WHITE, bold=True)
leverages = [
    "Existing PDU-Micro team provides power electronics and embedded firmware capability (3–4 engineers shared)",
    "EVerest open-source community provides protocol stack support, bug fixes, and conformance testing",
    "Standard industrial HVAC components — no custom refrigeration expertise needed full-time",
]
for l in leverages:
    add_bullet(tf, l, size=12, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 18: Risk Mitigation
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 18, "Risk Mitigation", "Top Technical Risks & Mitigations")

data = [
    ["Risk", "Impact", "Prob.", "Mitigation"],
    ["PDU-Micro firmware delayed", "Critical", "Med", "EVerest SIL simulation decouples EVerest work"],
    ["CAN-FD adapter incompatibility", "High", "Med", "Verify M1; PEAK PCAN-USB FD backup"],
    ["ISO 15118 PLC modem EV compat.", "High", "Med", "Test multiple EVs; DIN 70121 fallback"],
    ["Safety SIL 2 assessment fails", "High", "Low", "Engage assessor M3; design for SIL 2 upfront"],
    ["Thermal mgmt insufficient @150kW", "Med", "Low", "Thermal budget validated; full-power test M9"],
    ["OCPP 2.0.1 interop issues", "Med", "Med", "Multiple CSMS vendors; certified libocpp"],
    ["Current sharing instability (5 mod)", "Med", "Low", "Droop sharing (proven); early bench test"],
    ["SiC supply chain disruption", "Med", "Med", "Dual-source; 6-month buffer stock"],
]
table = make_table(slide, 9, 4, Inches(0.6), Inches(1.4), Inches(12.1), Inches(4.2))
style_table(table, data=data, col_widths=[Inches(3.0), Inches(1.2), Inches(0.9), Inches(7.0)])

# Color-code impact
impact_colors = {
    "Critical": ACCENT_RED,
    "High": ACCENT_ORANGE,
    "Med": RGBColor(0xCC, 0xCC, 0x00),
}
for row_idx in range(1, 9):
    cell = table.cell(row_idx, 1)
    text = cell.text
    if text in impact_colors:
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = impact_colors[text]
            p.font.bold = True

# Overall assessment
shape = add_shape(slide, Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.8), fill_color=RGBColor(0x0A, 0x1A, 0x30), border_color=ACCENT_GREEN)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(15)
tf.margin_top = Pt(8)
set_text(tf, "RISK ASSESSMENT:", size=12, color=ACCENT_GREEN, bold=True)
run = tf.paragraphs[0].add_run()
run.text = "  All critical risks have identified mitigations. Highest-risk item (PDU-Micro schedule) is mitigated by software simulation decoupling. No show-stoppers identified."
run.font.size = Pt(12)
run.font.color.rgb = LIGHT_GRAY
run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════
# SLIDE 19: Deployment Strategy
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, 19, "Deployment Strategy", "Three-Phase Market Rollout")

phases = [
    ("Phase 1: Pilot", "Months 14–18", ACCENT_BLUE, [
        "5–10 units at controlled sites",
        "Fleet depots, partner locations",
        "6-month field validation",
        "Collect MTBF, reliability data",
        "OTA firmware updates via OCPP",
    ]),
    ("Phase 2: Regional Launch", "Months 18–24", ACCENT_GREEN, [
        "50–100 units",
        "Highway corridors, urban hubs",
        "CPO partner integrations",
        "Dual-connector variant",
        "OCPP backend SaaS offering",
    ]),
    ("Phase 3: Scale Production", "Year 2+", ACCENT_ORANGE, [
        "500+ units/year",
        "300 kW and 350 kW variants",
        "International certifications",
        "In-house HVAC manufacturing",
        "Multi-connector power-sharing",
    ]),
]

for i, (title, timeline, color, items) in enumerate(phases):
    x = Inches(0.6) + i * Inches(4.15)
    card = add_shape(slide, x, Inches(1.4), Inches(3.95), Inches(4.0), fill_color=BG_CARD, border_color=color)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(12)
    set_text(tf, title, size=16, color=color, bold=True)
    add_paragraph(tf, timeline, size=11, color=MID_GRAY, space_before=Pt(2))
    for item in items:
        add_bullet(tf, item, size=12, color=LIGHT_GRAY)

# Revenue model
shape = add_shape(slide, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.0), fill_color=RGBColor(0x0A, 0x1A, 0x30), border_color=BORDER_COLOR)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(15)
tf.margin_top = Pt(10)
set_text(tf, "REVENUE MODEL", size=12, color=WHITE, bold=True)
add_paragraph(tf, "Hardware Sales  ·  OCPP Backend SaaS (optional)  ·  Maintenance Contracts  ·  Extended Warranty  ·  Power Module Upgrades", size=13, color=LIGHT_GRAY, space_before=Pt(6))


# ══════════════════════════════════════════════════════════════
# SLIDE 20: Investment Summary
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_GREEN
bar.line.fill.background()

add_title_bar(slide, 20, "Investment Summary", "Why Invest Now")

# Three columns
columns = [
    ("Technical Readiness", ACCENT_BLUE, [
        "30 kW SiC power module at Rev 0.7",
        "Core power electronics IP developed",
        "EVerest proven in production globally",
        "Safety designed to SIL 2 from inception",
        "Full documentation package complete",
    ]),
    ("Cost Advantage", ACCENT_GREEN, [
        "BOM: $21,500 @100 → $17,000 @1000",
        "$0 per-unit software licensing",
        "vs. $2K–$5K for competitors",
        "Target gross margin: 40–47%",
        "In-house design = supply chain control",
    ]),
    ("Market Position", ACCENT_ORANGE, [
        "Only open-source SW + in-house HW DCFC",
        "Highest efficiency in class (>96% SiC)",
        "Full 1000V — ready for 800V EV wave",
        "One platform: 50–350 kW",
        "No vendor lock-in for customers",
    ]),
]

for i, (title, color, items) in enumerate(columns):
    x = Inches(0.6) + i * Inches(4.15)
    card = add_shape(slide, x, Inches(1.4), Inches(3.95), Inches(3.5), fill_color=BG_CARD, border_color=color)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(12)
    set_text(tf, title, size=16, color=color, bold=True)
    for item in items:
        add_bullet(tf, item, size=12, color=LIGHT_GRAY)

# The ask
shape = add_shape(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.8), fill_color=RGBColor(0x0A, 0x1A, 0x30), border_color=ACCENT_GREEN)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(20)
tf.margin_top = Pt(12)
set_text(tf, "THE ASK", size=18, color=ACCENT_GREEN, bold=True)
asks = [
    "14-month development to certification-ready prototype",
    "8–10 person engineering team (3–4 shared from existing PDU-Micro team)",
    "Pilot deployment of 5–10 units at Month 14",
    "Path to 500+ units/year production by Year 2",
]
for a in asks:
    add_bullet(tf, a, size=14, color=WHITE)


# ── Save ──
output_path = "/Users/manaspradhan/Library/Mobile Documents/iCloud~md~obsidian/Documents/ClaudeNotes/__Workspaces/DCFC/DCFC Investor Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
